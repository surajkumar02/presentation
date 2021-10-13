from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
from wand.display import display
# from __future__ import print_function

ppt = Presentation()

images = ['image1.jpg', 'image2.jpg', 'image3.jpg', 'image4.jpg', 'image5.jpg']
logo ='nike_black.png'
size = Inches(4)
for i in range(0, len(images)):
    page = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(page)
    view = slide.shapes
    text = view.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
    text.text = f"Result page:{i+1}"
    text.text = 'Pic with Watermark'
    with Image(filename =images[i]) as image:
        with Image(filename =logo) as water:
            water.resize(800,350)
            image.watermark(water, 0.2, 50, 50)
            image.save(filename='new.jpg')
            view.add_picture('new.jpg', Inches(1), Inches(2), height=size)
ppt.save("result.pptx")
